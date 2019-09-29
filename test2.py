import os
from pptx import Presentation
import pypptxGenerator
import re

"""
pptx content detector should be in its own file
"""





# def find_text(prs):
# 	for slide in prs.slides:
# 	    for shape in slide.shapes:
# 	        if not shape.has_text_frame:
# 	        	print("Slide " + str((prs.slides.index(slide)) + 1) + " has NO text.")
# 	        	continue
# 	        for paragraph in shape.text_frame.paragraphs:
# 	            for run in paragraph.runs:
# 	            	print("Slide " + str((prs.slides.index(slide)) + 1) + " has text:")
# 	            	# text_runs.append(run.text)
# 	            	print(run.text)


def find_table(prs, __file__):
	for slide in prs.slides:
		for shape in slide.shapes:
			try:
				if shape.table:
					print("Slide " + str((prs.slides.index(slide)) + 1) + " has table.")
					#writeTable(__file__)
					continue
			except:
				# print("Slide " + str((prs.slides.index(slide)) + 1) + " has NO table.")
				pass
			# for paragraph in shape.text_frame.paragraphs:
			#     for run in paragraph.runs:
			#     	print("Slide " + str((prs.slides.index(slide)) + 1) + " has text:")
			#     	# text_runs.append(run.text)
			#     	print(run.text)

def find_text(prs):
	for slide in prs.slides:
		for shape in slide.shapes:
			try:
				if shape.text_frame:
					print("Slide " + str((prs.slides.index(slide)) + 1) + " has text_frame.")
					for paragraph in shape.text_frame.paragraphs:
						text_runs = []
						for run in paragraph.runs:
							text_runs.append(run.text)
							print(run.text)
						print("Slide " + str((prs.slides.index(slide)) + 1) + " has text paragraph: \n" + test_runs)
					#writeTable(__file__)
					continue
			except:
				# print("Slide " + str((prs.slides.index(slide)) + 1) + " has NO table.")
				pass

def string_found(string1, string2):
   if re.search(r"\b" + re.escape(string1) + r"\b", string2):
      return True
   return False

def find_placeholder(prs):
	for slide in prs.slides:
		for shape in slide.placeholders:
			try:
				if shape.placeholder_format:
					print("\n\n#############################################################\n\n")
					print("Slide " + str((prs.slides.index(slide)) + 1) + " has placeholder_format.")
					print("\nShape Type: " + str(shape.placeholder_format.type))

					"""
					USE ALL CAPS
					BITMAP
					    Bitmap
					BODY
					    Body
					CENTER_TITLE
					    Center Title
					CHART
					    Chart
					DATE
					    Date
					FOOTER
					    Footer
					HEADER
					    Header
					MEDIA_CLIP
					    Media Clip
					OBJECT
					    Object
					ORG_CHART
					    Organization Chart
					PICTURE
					    Picture
					SLIDE_NUMBER
					    Slide Number
					SUBTITLE
					    Subtitle
					TABLE
					    Table
					TITLE
					    Title
					VERTICAL_BODY
					    Vertical Body
					VERTICAL_OBJECT
					    Vertical Object
					VERTICAL_TITLE
					    Vertical Title
					MIXED

					"""
					if(string_found('TITLE',str(shape.placeholder_format.type))):
						print("\n FOUND A TITLE")
						if(shape.text_frame):
							print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
						else:
							print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
					if(string_found('TABLE',str(shape.placeholder_format.type))):
						print("\n FOUND A TABLE")
					if(string_found('PICTURE',str(shape.placeholder_format.type))):
						print("\n FOUND A PICTURE")
					if(string_found('CHART',str(shape.placeholder_format.type))):
						print("\n FOUND A CHART")
					if(string_found('BODY',str(shape.placeholder_format.type))):
						print("\n FOUND A BODY")
						if(shape.text_frame):
							print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
						elif(shape.text):
							print("\nSHAPE.TEXT: " + shape.text)
						else:
							print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
					if(string_found('OBJECT',str(shape.placeholder_format.type))):
						print("\n FOUND A OBJECT")
						if(shape.text_frame):
							print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
						else:
							print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
					print('\tidx: %d \n\tname: %s\n\tplaceholder_format_type: %s\n\tshape.shape_type: %s' % (shape.placeholder_format.idx, shape.name, str(shape.placeholder_format.type), shape.shape_type))
					continue
			except Exception as detail:
				print(detail)
				pass



if __name__ == "__main__":
	## generate file "mergeMe.py"
	# pptx2py()	    		
	current_dr = os.path.dirname(os.path.realpath(__file__))
	prs = Presentation(current_dr + '/template_sample.pptx')
	find_placeholder(prs)