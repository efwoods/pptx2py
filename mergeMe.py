from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("./.template/DO_NOT_CHANGE.pptx")
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 685800
TOP = 2895480
WIDTH = 7771320
HEIGHT = 1065600
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Presentation Title'
shape.text_frame.word_wrap = None
LEFT = 685800
TOP = 4648320
WIDTH = 6399720
HEIGHT = 1904040
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Author,Department,Date,Location'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 1086840
TOP = 496440
WIDTH = 7771320
HEIGHT = 1361160
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Presentation Title'
shape.text_frame.word_wrap = None
LEFT = 274320
TOP = 2377440
WIDTH = 7771320
HEIGHT = 1682640
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Author,Department,Date,Location'
shape.text_frame.word_wrap = None
LEFT = 1645920
TOP = 1463040
WIDTH = 7040520
HEIGHT = 4937400
ROWS = 4
COLS = 9
shape = slide.shapes.add_table(ROWS, COLS, LEFT, TOP, WIDTH, HEIGHT)
for row in range (0, len(shape.table.rows)):
	for col in range (0,len(shape.table.columns)):
		if(row == 0):
			shape.table.columns[0].width = 781560
			shape.table.cell(0,0).text = ''
			shape.table.cell(0,0).fill.solid()
			shape.table.cell(0,0).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[1].width = 781560
			shape.table.cell(0,1).text = ''
			shape.table.cell(0,1).fill.solid()
			shape.table.cell(0,1).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[2].width = 781560
			shape.table.cell(0,2).text = ''
			shape.table.cell(0,2).fill.solid()
			shape.table.cell(0,2).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[3].width = 781560
			shape.table.cell(0,3).text = ''
			shape.table.cell(0,3).fill.solid()
			shape.table.cell(0,3).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[4].width = 781560
			shape.table.cell(0,4).text = ''
			shape.table.cell(0,4).fill.solid()
			shape.table.cell(0,4).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[5].width = 781560
			shape.table.cell(0,5).text = ''
			shape.table.cell(0,5).fill.solid()
			shape.table.cell(0,5).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[6].width = 781560
			shape.table.cell(0,6).text = ''
			shape.table.cell(0,6).fill.solid()
			shape.table.cell(0,6).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[7].width = 781560
			shape.table.cell(0,7).text = ''
			shape.table.cell(0,7).fill.solid()
			shape.table.cell(0,7).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[8].width = 788400
			shape.table.cell(0,8).text = ''
			shape.table.cell(0,8).fill.solid()
			shape.table.cell(0,8).fill.fore_color.rgb = RGBColor(179, 179, 179)


		shape.table.cell(1,0).text = ''
		shape.table.cell(1,0).fill.solid()
		shape.table.cell(1,0).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,1).text = ''
		shape.table.cell(1,1).fill.solid()
		shape.table.cell(1,1).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,2).text = ''
		shape.table.cell(1,2).fill.solid()
		shape.table.cell(1,2).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,3).text = ''
		shape.table.cell(1,3).fill.solid()
		shape.table.cell(1,3).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,4).text = ''
		shape.table.cell(1,4).fill.solid()
		shape.table.cell(1,4).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,5).text = ''
		shape.table.cell(1,5).fill.solid()
		shape.table.cell(1,5).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,6).text = ''
		shape.table.cell(1,6).fill.solid()
		shape.table.cell(1,6).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,7).text = ''
		shape.table.cell(1,7).fill.solid()
		shape.table.cell(1,7).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(1,8).text = ''
		shape.table.cell(1,8).fill.solid()
		shape.table.cell(1,8).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(2,0).text = 'Great new updates =)'
		shape.table.cell(2,0).fill.solid()
		shape.table.cell(2,0).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,1).text = ''
		shape.table.cell(2,1).fill.solid()
		shape.table.cell(2,1).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,2).text = ''
		shape.table.cell(2,2).fill.solid()
		shape.table.cell(2,2).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,3).text = ''
		shape.table.cell(2,3).fill.solid()
		shape.table.cell(2,3).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,4).text = ''
		shape.table.cell(2,4).fill.solid()
		shape.table.cell(2,4).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,5).text = ''
		shape.table.cell(2,5).fill.solid()
		shape.table.cell(2,5).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,6).text = ''
		shape.table.cell(2,6).fill.solid()
		shape.table.cell(2,6).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,7).text = ''
		shape.table.cell(2,7).fill.solid()
		shape.table.cell(2,7).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(2,8).text = ''
		shape.table.cell(2,8).fill.solid()
		shape.table.cell(2,8).fill.fore_color.rgb = RGBColor(230, 230, 230)


		shape.table.cell(3,0).text = ''
		shape.table.cell(3,0).fill.solid()
		shape.table.cell(3,0).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,1).text = ''
		shape.table.cell(3,1).fill.solid()
		shape.table.cell(3,1).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,2).text = ''
		shape.table.cell(3,2).fill.solid()
		shape.table.cell(3,2).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,3).text = ''
		shape.table.cell(3,3).fill.solid()
		shape.table.cell(3,3).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,4).text = ''
		shape.table.cell(3,4).fill.solid()
		shape.table.cell(3,4).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,5).text = ''
		shape.table.cell(3,5).fill.solid()
		shape.table.cell(3,5).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,6).text = ''
		shape.table.cell(3,6).fill.solid()
		shape.table.cell(3,6).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,7).text = ''
		shape.table.cell(3,7).fill.solid()
		shape.table.cell(3,7).fill.fore_color.rgb = RGBColor(204, 204, 204)


		shape.table.cell(3,8).text = ''
		shape.table.cell(3,8).fill.solid()
		shape.table.cell(3,8).fill.fore_color.rgb = RGBColor(204, 204, 204)


slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 722160
TOP = 4648320
WIDTH = 7771320
HEIGHT = 1682640
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = ''
shape.text_frame.word_wrap = None
LEFT = 731520
TOP = 1188720
WIDTH = 7589160
HEIGHT = 5394600
ROWS = 2
COLS = 5
shape = slide.shapes.add_table(ROWS, COLS, LEFT, TOP, WIDTH, HEIGHT)
for row in range (0, len(shape.table.rows)):
	for col in range (0,len(shape.table.columns)):
		if(row == 0):
			shape.table.columns[0].width = 1517040
			shape.table.cell(0,0).text = 'new'
			shape.table.cell(0,0).fill.solid()
			shape.table.cell(0,0).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[1].width = 1517040
			shape.table.cell(0,1).text = 'table'
			shape.table.cell(0,1).fill.solid()
			shape.table.cell(0,1).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[2].width = 1517040
			shape.table.cell(0,2).text = 'here'
			shape.table.cell(0,2).fill.solid()
			shape.table.cell(0,2).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[3].width = 1517040
			shape.table.cell(0,3).text = 'this'
			shape.table.cell(0,3).fill.solid()
			shape.table.cell(0,3).fill.fore_color.rgb = RGBColor(179, 179, 179)


			shape.table.columns[4].width = 1521360
			shape.table.cell(0,4).text = 'works'
			shape.table.cell(0,4).fill.solid()
			shape.table.cell(0,4).fill.fore_color.rgb = RGBColor(179, 179, 179)


		shape.table.cell(1,0).text = 'Alex, Evan'
		shape.table.cell(1,0).fill.solid()
		shape.table.cell(1,0).fill.fore_color.rgb = RGBColor(173, 213, 138)


		shape.table.cell(1,1).text = '10/02/2019'
		shape.table.cell(1,1).fill.solid()
		shape.table.cell(1,1).fill.fore_color.rgb = RGBColor(173, 213, 138)


		shape.table.cell(1,2).text = 'Making new '
		shape.table.cell(1,2).fill.solid()
		shape.table.cell(1,2).fill.fore_color.rgb = RGBColor(173, 213, 138)


		shape.table.cell(1,3).text = 'changes'
		shape.table.cell(1,3).fill.solid()
		shape.table.cell(1,3).fill.fore_color.rgb = RGBColor(173, 213, 138)


		shape.table.cell(1,4).text = 'everyday'
		shape.table.cell(1,4).fill.solid()
		shape.table.cell(1,4).fill.fore_color.rgb = RGBColor(173, 213, 138)


LEFT = 2834640
TOP = 274320
WIDTH = 3474000
HEIGHT = 601560
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Title: there was a black text box…?'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 533520
TOP = 990720
WIDTH = 3007080
HEIGHT = 551520
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'GraphTitle'
shape.text_frame.word_wrap = None
LEFT = 533520
TOP = 6019920
WIDTH = 8228520
HEIGHT = 684720
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Additional Notes ,- Lorem ipsum dolor sit amet, consectetuer adipiscing elit. Aenean commodo ligula eget dolor. Aenean massa. Cum sociis natoque penatibus et magnis dis parturient montes, nascetur ridiculus mus. Donec quam felis, ultricies nec, pellentesque eu, pretium quis, sem'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 533520
TOP = 990720
WIDTH = 8152200
HEIGHT = 532440
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Chart Title'
shape.text_frame.word_wrap = None
LEFT = 529560
TOP = 1667520
WIDTH = 8152920
HEIGHT = 4419000
ROWS = 8
COLS = 4
shape = slide.shapes.add_table(ROWS, COLS, LEFT, TOP, WIDTH, HEIGHT)
for row in range (0, len(shape.table.rows)):
	for col in range (0,len(shape.table.columns)):
		if(row == 0):
			shape.table.columns[0].width = 2038320
			shape.table.cell(0,0).text = 'Column A'
			shape.table.cell(0,0).fill.solid()
			shape.table.cell(0,0).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[1].width = 2038320
			shape.table.cell(0,1).text = 'B'
			shape.table.cell(0,1).fill.solid()
			shape.table.cell(0,1).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[2].width = 2038320
			shape.table.cell(0,2).text = 'C'
			shape.table.cell(0,2).fill.solid()
			shape.table.cell(0,2).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[3].width = 2038320
			shape.table.cell(0,3).text = 'D'
			shape.table.cell(0,3).fill.solid()
			shape.table.cell(0,3).fill.fore_color.rgb = RGBColor(192, 0, 0)


		shape.table.cell(1,0).text = 'XXXXXXXX'
		shape.table.cell(1,0).fill.solid()
		shape.table.cell(1,0).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(1,1).text = 'XX'
		shape.table.cell(1,1).fill.solid()
		shape.table.cell(1,1).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(1,2).text = 'XX'
		shape.table.cell(1,2).fill.solid()
		shape.table.cell(1,2).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(1,3).text = 'XX'
		shape.table.cell(1,3).fill.solid()
		shape.table.cell(1,3).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(2,0).text = 'XXXXXXXX'
		shape.table.cell(2,0).fill.solid()
		shape.table.cell(2,0).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(2,1).text = 'XX'
		shape.table.cell(2,1).fill.solid()
		shape.table.cell(2,1).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(2,2).text = 'XX'
		shape.table.cell(2,2).fill.solid()
		shape.table.cell(2,2).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(2,3).text = 'XX'
		shape.table.cell(2,3).fill.solid()
		shape.table.cell(2,3).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(3,0).text = 'XXXXXXXX'
		shape.table.cell(3,0).fill.solid()
		shape.table.cell(3,0).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(3,1).text = 'XX'
		shape.table.cell(3,1).fill.solid()
		shape.table.cell(3,1).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(3,2).text = 'XX'
		shape.table.cell(3,2).fill.solid()
		shape.table.cell(3,2).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(3,3).text = 'XX'
		shape.table.cell(3,3).fill.solid()
		shape.table.cell(3,3).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(4,0).text = 'XXXXXXXX'
		shape.table.cell(4,0).fill.solid()
		shape.table.cell(4,0).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(4,1).text = 'XX'
		shape.table.cell(4,1).fill.solid()
		shape.table.cell(4,1).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(4,2).text = 'XX'
		shape.table.cell(4,2).fill.solid()
		shape.table.cell(4,2).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(4,3).text = 'XX'
		shape.table.cell(4,3).fill.solid()
		shape.table.cell(4,3).fill.fore_color.rgb = RGBColor(244, 231, 231)


		shape.table.cell(5,0).text = 'XXXXXXXX'
		shape.table.cell(5,0).fill.solid()
		shape.table.cell(5,0).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(5,1).text = 'XX'
		shape.table.cell(5,1).fill.solid()
		shape.table.cell(5,1).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(5,2).text = 'XX'
		shape.table.cell(5,2).fill.solid()
		shape.table.cell(5,2).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(5,3).text = 'XX'
		shape.table.cell(5,3).fill.solid()
		shape.table.cell(5,3).fill.fore_color.rgb = RGBColor(232, 204, 204)


		shape.table.cell(6,0).text = 'XXXXXXXX'
		shape.table.cell(6,0).fill.solid()
		shape.table.cell(6,0).fill.fore_color.rgb = RGBColor(69, 79, 161)


		shape.table.cell(6,1).text = 'Oct 19'
		shape.table.cell(6,1).fill.solid()
		shape.table.cell(6,1).fill.fore_color.rgb = RGBColor(94, 138, 199)


		shape.table.cell(6,2).text = 'Updated and complete'
		shape.table.cell(6,2).fill.solid()
		shape.table.cell(6,2).fill.fore_color.rgb = RGBColor(94, 138, 199)


		shape.table.cell(6,3).text = 'hurrah'
		shape.table.cell(6,3).fill.solid()
		shape.table.cell(6,3).fill.fore_color.rgb = RGBColor(27, 117, 188)


		shape.table.cell(7,0).text = 'Alex, Evan'
		shape.table.cell(7,0).fill.solid()
		shape.table.cell(7,0).fill.fore_color.rgb = RGBColor(140, 207, 183)


		shape.table.cell(7,1).text = 'Oct 23423'
		shape.table.cell(7,1).fill.solid()
		shape.table.cell(7,1).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(7,2).text = 'New assignments'
		shape.table.cell(7,2).fill.solid()
		shape.table.cell(7,2).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(7,3).text = '_)'
		shape.table.cell(7,3).fill.solid()
		shape.table.cell(7,3).fill.fore_color.rgb = RGBColor(0, 178, 116)


slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 2468880
TOP = 457200
WIDTH = 4023360
HEIGHT = 346320
shape = slide.shapes.add_textbox(LEFT,TOP,WIDTH,HEIGHT)
shape.line.dash_style = None
shape.line.width = 0
LEFT = 538200
TOP = 1109160
WIDTH = 8152920
HEIGHT = 4419000
ROWS = 3
COLS = 4
shape = slide.shapes.add_table(ROWS, COLS, LEFT, TOP, WIDTH, HEIGHT)
for row in range (0, len(shape.table.rows)):
	for col in range (0,len(shape.table.columns)):
		if(row == 0):
			shape.table.columns[0].width = 2038320
			shape.table.cell(0,0).text = 'Column A'
			shape.table.cell(0,0).fill.solid()
			shape.table.cell(0,0).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[1].width = 2038320
			shape.table.cell(0,1).text = 'B'
			shape.table.cell(0,1).fill.solid()
			shape.table.cell(0,1).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[2].width = 2038320
			shape.table.cell(0,2).text = 'C'
			shape.table.cell(0,2).fill.solid()
			shape.table.cell(0,2).fill.fore_color.rgb = RGBColor(192, 0, 0)


			shape.table.columns[3].width = 2038320
			shape.table.cell(0,3).text = 'D'
			shape.table.cell(0,3).fill.solid()
			shape.table.cell(0,3).fill.fore_color.rgb = RGBColor(192, 0, 0)


		shape.table.cell(1,0).text = 'Alex, Evan'
		shape.table.cell(1,0).fill.solid()
		shape.table.cell(1,0).fill.fore_color.rgb = RGBColor(140, 207, 183)


		shape.table.cell(1,1).text = 'Oct 23423'
		shape.table.cell(1,1).fill.solid()
		shape.table.cell(1,1).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(1,2).text = 'New assignments'
		shape.table.cell(1,2).fill.solid()
		shape.table.cell(1,2).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(1,3).text = '_)'
		shape.table.cell(1,3).fill.solid()
		shape.table.cell(1,3).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(2,0).text = 'people'
		shape.table.cell(2,0).fill.solid()
		shape.table.cell(2,0).fill.fore_color.rgb = RGBColor(140, 207, 183)


		shape.table.cell(2,1).text = 'As;dlfkja;lsdkj'
		shape.table.cell(2,1).fill.solid()
		shape.table.cell(2,1).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(2,2).text = 'More updates'
		shape.table.cell(2,2).fill.solid()
		shape.table.cell(2,2).fill.fore_color.rgb = RGBColor(0, 178, 116)


		shape.table.cell(2,3).text = 'asd'
		shape.table.cell(2,3).fill.solid()
		shape.table.cell(2,3).fill.fore_color.rgb = RGBColor(0, 178, 116)


slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 533520
TOP = 990720
WIDTH = 8152200
HEIGHT = 532440
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'List Title'
shape.text_frame.word_wrap = None
LEFT = 533520
TOP = 1676520
WIDTH = 8152200
HEIGHT = 4448520
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Lorem ipsum dolor sit amet ,Aenean commodo ligula eget dolor ,Cum sociis natoque penatibus et magnis dis parturient montes ,Donec quam felis, ultricies nec, pellentesque eu ,Lorem ipsum dolor sit amet, consectetuer adipiscing elit ,Aenean massa ,Aenean commodo ligula eget dolor'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 533520
TOP = 990720
WIDTH = 8152200
HEIGHT = 532440
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Slide Title'
shape.text_frame.word_wrap = None
LEFT = 533520
TOP = 1676520
WIDTH = 8152200
HEIGHT = 4448520
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Lorem ipsum dolor sit amet, consectetuer adipiscing elit. Aenean commodo ligula eget dolor. Aenean massa. Cum sociis natoque penatibus et magnis dis parturient montes, nascetur ridiculus mus. Donec quam felis, ultricies nec, pellentesque eu, pretium quis, sem.,Aenean massa. Cum sociis natoque penatibus et magnis dis parturient montes, nascetur ridiculus mus. Donec quam felis, ultricies nec, pellentesque eu, pretium quis, sem.'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 533520
TOP = 6248520
WIDTH = 8152200
HEIGHT = 346680
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Type image caption here.'
shape.text_frame.word_wrap = None
LEFT = 533520
TOP = 1066680
WIDTH = 8152200
HEIGHT = 456120
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Slide Title'
shape.text_frame.word_wrap = None
slide = prs.slides.add_slide(title_slide_layout)
for shape in slide.shapes:
	sp = shape.element
	sp.getparent().remove(sp)
LEFT = 507960
TOP = 1038240
WIDTH = 8228520
HEIGHT = 531720
shape = slide.shapes.add_shape(
	MSO_SHAPE.RECTANGLE, LEFT, TOP, WIDTH, HEIGHT)
current_paragraph = shape.text_frame.add_paragraph()
current_paragraph.font.bold = None
current_paragraph.font.name = None
current_paragraph.font.size = None
current_paragraph.font.underline = None
current_paragraph.font.italic = None
shape.text_frame.text = 'Page Title'
shape.text_frame.word_wrap = None
prs.save('FINAL.pptx')
