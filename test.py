import os
from pptx import Presentation

current_dr = os.path.dirname(os.path.realpath(__file__))
prs = Presentation(current_dr + '/template_sample.pptx')
# newPrs = Presentation()

numberOfSlides = len(prs.slides)
if(numberOfSlides > 0):
	print("I have counted " + str(numberOfSlides) + " slides.")
	print("Checking for slide notes...")
	i = 0
	while (i < numberOfSlides):
		if(prs.slides[i].has_notes_slide):
			print("Slide " + str(i + 1) + " has notes.")
			print(prs.slides[i].notes_slide.notes_placeholder.text)
			print('\n')
			print('\n')
			print('\n')
		# 	if(prs.slides[i].notes_slide.notes_placeholder.has_chart):
		# 		print("Slide " + str(i + 1) + " has chart.")

		# # print("chart")
		# if(print(prs.slides[i].notes_slide.notes_placeholder.has_chart)):
		# 	print("chart")
		
		# # print("table")
		# if(print(prs.slides[i].notes_slide.notes_placeholder.has_table)):
		# 	print("table")
		
		i = i + 1
else:
	print("I did not see any slides to count.")



# i = 5
# # while ( i < 2):
# title_slide_layout = prs.slide_layouts[i]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# # subtitle = slide.placeholders[1]
# # i = i + 1


# # title.text = "Hello, World!"
# # subtitle.text = "python-pptx was here!"

# prs.save('test-template-update.pptx')

