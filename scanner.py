"""
############################################################################
IMPORTS
############################################################################
"""
import types
import os
from pptx import Presentation
import re
import sys
from pprint import pprint

"""
############################################################################
UTILITY FUNCTIONS
############################################################################
"""

def string_found(string1, string2):
   if re.search(r"\b" + re.escape(string1) + r"\b", string2):
      return True
   return False

def findShapeDimensions(shape):
	print("\n SHAPE LEFT: " + str(shape.left))
	print("\n SHAPE TOP: " + str(shape.top))
	print("\n SHAPE WIDTH: " + str(shape.width))
	print("\n SHAPE HEIGHT: " + str(shape.height))

def separateFoundShapesWithPrint(FOUND_BY):
	print('\n\
		XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX\
		\n\
		...FOUND SHAPE: EXITING ID SHAPE ' + FOUND_BY + ' ...\
		\n\
		XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXX\
		\n')

"""
############################################################################
FIND BY SHAPE_TYPE FUNCTIONS
############################################################################
"""



def printAttrTree(shape):
	for attr in dir(shape):
				if not attr.startswith('__'):
					try:
						print(attr + ":" + "\n")
						
						if not (type(getattr(shape,attr)) == str):
							#printAttrTree(attr)
							print(getattr(shape,attr))
						else:
							print(getattr(shape, attr))
						
						print("\n")
						#printAttrTree(attr)
					except Exception as detail:
						print(detail)
						continue

def findShapeAUTO_SHAPE(shape):
	if(string_found('AUTO_SHAPE',str(shape.shape_type))):
		print("\n FOUND A AUTO_SHAPE")
		findShapeDimensions(shape)
		try:
			print("\n")

			#printAttrTree(shape)
			""" 
			print('shape.table.cell(0,0).fill.type: ', shape..fill.type)
		

			print("\nTABLE CELL text_frame: ")
			pprint(dir(shape.table.cell(0,0).text_frame))
			print("\nTABLE CELL text_frame.text: ")
			print(shape.table.cell(0,0).text_frame.text)
			
			print("\nCell FONT PER PARAGRAPH")
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.bold: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.bold)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.color: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.color.rgb)
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font.color))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.fill: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.fill)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.italic: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.italic)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.name: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.name)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.size: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.size)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.underline: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.underline)
 """
		except Exception as detail:
			print(detail)
			pass

		return True
	return False

def findShapeCALLOUT(shape):
	if(string_found('CALLOUT',str(shape.shape_type))):
		print("\n FOUND A CALLOUT")
		return True
	return False

def findShapeCANVAS(shape):
	if(string_found('CANVAS',str(shape.shape_type))):
		print("\n FOUND A CANVAS")
		return True
	return False

def findShapeCHART(shape):
	if(string_found('CHART',str(shape.shape_type))):
		print("\n FOUND A CHART")
		findShapeDimensions(shape)
		return True
	return False

def findShapeCOMMENT(shape):
	if(string_found('COMMENT',str(shape.shape_type))):
		print("\n FOUND A COMMENT")
		return True
	return False

def findShapeDIAGRAM(shape):
	if(string_found('DIAGRAM',str(shape.shape_type))):
		print("\n FOUND A DIAGRAM")
		return True
	return False

def findShapeEMBEDDED_OLE_OBJECT(shape):
	if(string_found('EMBEDDED_OLE_OBJECT',str(shape.shape_type))):
		print("\n FOUND A EMBEDDED_OLE_OBJECT")
		return True
	return False

def findShapeFORM_CONTROL(shape):
	if(string_found('FORM_CONTROL',str(shape.shape_type))):
		print("\n FOUND A FORM_CONTROL")
		return True
	return False

def findShapeFREEFORM(shape):
	if(string_found('FREEFORM',str(shape.shape_type))):
		print("\n FOUND A FREEFORM")
		return True
	return False

def findShapeGROUP(shape):
	if(string_found('GROUP',str(shape.shape_type))):
		print("\n FOUND A GROUP")
		return True
	return False

def findShapeIGX_GRAPHIC(shape):
	if(string_found('IGX_GRAPHIC',str(shape.shape_type))):
		print("\n FOUND A IGX_GRAPHIC")
		return True
	return False

def findShapeINK(shape):
	if(string_found('INK',str(shape.shape_type))):
		print("\n FOUND A INK")
		return True
	return False

def findShapeINK_COMMENT(shape):
	if(string_found('INK_COMMENT',str(shape.shape_type))):
		print("\n FOUND A INK_COMMENT")
		return True
	return False

def findShapeLINE(shape):
	if(string_found('LINE',str(shape.shape_type))):
		print("\n FOUND A LINE")
		return True
	return False

def findShapeLINKED_OLE_OBJECT(shape):
	if(string_found('LINKED_OLE_OBJECT',str(shape.shape_type))):
		print("\n FOUND A LINKED_OLE_OBJECT")
		return True
	return False

def findShapeLINKED_PICTURE(shape):
	if(string_found('LINKED_PICTURE',str(shape.shape_type))):
		print("\n FOUND A LINKED_PICTURE")
		return True
	return False

def findShapeMEDIA(shape):
	if(string_found('MEDIA',str(shape.shape_type))):
		print("\n FOUND A MEDIA")
		return True
	return False

def findShapeOLE_CONTROL_OBJECT(shape):
	if(string_found('OLE_CONTROL_OBJECT',str(shape.shape_type))):
		print("\n FOUND A OLE_CONTROL_OBJECT")
		return True
	return False

def findShapePICTURE(shape):
	if(string_found('PICTURE',str(shape.shape_type))):
		print("\n FOUND A PICTURE")
		findShapeDimensions(shape)
		return True
	return False

def findShapePLACEHOLDER(shape):
	if(string_found('PLACEHOLDER',str(shape.shape_type))):
		print("\n FOUND A PLACEHOLDER")
		return True
	return False

def findShapeSCRIPT_ANCHOR(shape):
	if(string_found('SCRIPT_ANCHOR',str(shape.shape_type))):
		print("\n FOUND A SCRIPT_ANCHOR")
		return True
	return False

def findShapeTABLE_ATTRIBUTES(shape):
	if(string_found('TABLE',str(shape.shape_type))):
		print("\n FOUND A TABLE")
		findShapeDimensions(shape)
		
		print("\nShape Attributes: ")
		pprint(dir(shape))
		print("\nTable Attributes")
		pprint(dir(shape.table))
		print("\nTABLE ROWS: " 
		# number of rows in table
		+ str(len(shape.table.rows)))
		pprint(dir(shape.table.rows))
		print("\nTABLE COLUMNS: "
		# number of columns in table
		+ str(len(shape.table.columns)))
		pprint(dir(shape.table.columns))
		print("\nTABLE shape.table.cell")
		pprint(dir(shape.table.cell))
		print("\n")
		print("\nTABLE CELL(0,0)")
		pprint(dir(shape.table.cell(0,0)))
		print("\n")
		print("\nTABLE CELL(0,0) text: ")
		print(shape.table.cell(0,0).text)
		print("\nTABLE CELL FILL: ")
		pprint(dir(shape.table.cell(0,0).fill))

		""" try:
			pprint('shape.table.cell(0,0).fill.back_color: ', dir(shape.table.cell(0,0).fill.back_color))
		except Exception as detail:
			print(detail)
			pass 
		"""
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.background: ' + str(shape.table.cell(0,0).fill.background()))
		except Exception as detail:
			print(detail)
			pass
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.fore_color: ')
			pprint(dir(shape.table.cell(0,0).fill.fore_color))
			print("\nTABLE CELL COLOR")
			print(str(shape.table.cell(0,0).fill.fore_color.rgb))
		except Exception as detail:
			print(detail)
			pass
		"""

		try:
			print("\n")
			print('shape.table.cell(0,0).fill.from_fill_parent: ', shape.table.cell(0,0).fill.from_fill_parent)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient: ', shape.table.cell(0,0).fill.gradient)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient_angle: ', shape.table.cell(0,0).fill.gradient_angle)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient_stops: ', shape.table.cell(0,0).fill.gradient_stops)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.pattern: ', shape.table.cell(0,0).fill.pattern)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.patterned: ', shape.table.cell(0,0).fill.patterned)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.solid: ', shape.table.cell(0,0).fill.solid)
		except Exception as detail:
			print(detail)
			pass
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.type: ', shape.table.cell(0,0).fill.type)
		except Exception as detail:
			print(detail)
			pass

		print("\nTABLE CELL text_frame: ")
		pprint(dir(shape.table.cell(0,0).text_frame))
		print("\nTABLE CELL text_frame.text: ")
		print(shape.table.cell(0,0).text_frame.text)
		try:
			print("\nCell FONT PER PARAGRAPH")
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.bold: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.bold)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.color: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.color.rgb)
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font.color))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.fill: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.fill)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.italic: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.italic)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.name: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.name)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.size: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.size)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.underline: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.underline)
		except Exception as detail:
			print(detail)
			pass

		print("\n")

		print("\nTABLE COLUMN HEADERS & CONTENT")
		for row in range (0, len(shape.table.rows)):
			for col in range (0,len(shape.table.columns)):
				# cell (rows, columns)
				if(row == 0):
					print("Column #" + str(col + 1) + " Header: " + shape.table.cell(row,col).text + "\n")
				else:
					print("Cell(ROW, COL): (" +str(row + 1) + ", " + str(col + 1) + ") TEXT: " + shape.table.cell(row,col).text + "\n")

		""" for col in shape.table.columns:
			print("\nCOLUMN %s WIDTH: %s") % str(col), str(col.width)
		print("\n") """
		""" for i in range(0, len(shape.table.columns)):
			print("\nCOLUMN %d HEADER: %s") % i, str(shape.table.cell(0, i).text) 
		"""
		print("\n")


		return True
	return False

def findShapeTABLE(shape):
	if(string_found('TABLE',str(shape.shape_type))):
		print("\n FOUND A TABLE")
		findShapeDimensions(shape)
		
		print("\nShape Attributes: ")
		pprint(dir(shape))
		print("\nTable Attributes")
		pprint(dir(shape.table))
		print("\nTABLE ROWS: " 
		# number of rows in table
		+ str(len(shape.table.rows)))
		pprint(dir(shape.table.rows))
		print("\nTABLE COLUMNS: "
		# number of columns in table
		+ str(len(shape.table.columns)))
		pprint(dir(shape.table.columns))
		print("\nTABLE shape.table.cell")
		pprint(dir(shape.table.cell))
		print("\n")
		print("\nTABLE CELL(0,0)")
		pprint(dir(shape.table.cell(0,0)))
		print("\n")
		print("\nTABLE CELL(0,0) text: ")
		print(shape.table.cell(0,0).text)
		print("\nTABLE CELL FILL: ")
		pprint(dir(shape.table.cell(0,0).fill))

		""" try:
			pprint('shape.table.cell(0,0).fill.back_color: ', dir(shape.table.cell(0,0).fill.back_color))
		except Exception as detail:
			print(detail)
			pass 
		"""
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.background: ' + str(shape.table.cell(0,0).fill.background()))
		except Exception as detail:
			print(detail)
			pass
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.fore_color: ')
			pprint(dir(shape.table.cell(0,0).fill.fore_color))
			print("\nTABLE CELL COLOR")
			print(str(shape.table.cell(0,0).fill.fore_color.rgb))
		except Exception as detail:
			print(detail)
			pass
		"""

		try:
			print("\n")
			print('shape.table.cell(0,0).fill.from_fill_parent: ', shape.table.cell(0,0).fill.from_fill_parent)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient: ', shape.table.cell(0,0).fill.gradient)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient_angle: ', shape.table.cell(0,0).fill.gradient_angle)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.gradient_stops: ', shape.table.cell(0,0).fill.gradient_stops)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.pattern: ', shape.table.cell(0,0).fill.pattern)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.patterned: ', shape.table.cell(0,0).fill.patterned)
		except Exception as detail:
			print(detail)
			pass
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.solid: ', shape.table.cell(0,0).fill.solid)
		except Exception as detail:
			print(detail)
			pass
		"""
		try:
			print("\n")
			print('shape.table.cell(0,0).fill.type: ', shape.table.cell(0,0).fill.type)
		

			print("\nTABLE CELL text_frame: ")
			pprint(dir(shape.table.cell(0,0).text_frame))
			print("\nTABLE CELL text_frame.text: ")
			print(shape.table.cell(0,0).text_frame.text)
			
			print("\nCell FONT PER PARAGRAPH")
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.bold: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.bold)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.color: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.color.rgb)
			pprint(dir(shape.table.cell(0,0).text_frame.paragraphs[0].font.color))

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.fill: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.fill)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.italic: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.italic)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.language_id)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.name: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.name)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.size: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.size)

			print('shape.table.cell(0,0).text_frame.paragraphs[0].font.underline: ', shape.table.cell(0,0).text_frame.paragraphs[0].font.underline)

		except Exception as detail:
			print(detail)
			pass

		print("\n")
		## TABLE LINES 
		print("\nTABLE COLUMN HEADERS & CONTENT")
		for row in range (0, len(shape.table.rows)):
			for col in range (0,len(shape.table.columns)):
				# cell (rows, columns)
				if(row == 0):
					print("Column #" + str(col + 1) + " Header: " + shape.table.cell(row,col).text + "\n")
				else:
					print("Cell(ROW, COL): (" +str(row + 1) + ", " + str(col + 1) + ") TEXT: " + shape.table.cell(row,col).text + "\n")

		""" for col in shape.table.columns:
			print("\nCOLUMN %s WIDTH: %s") % str(col), str(col.width)
		print("\n") """
		""" for i in range(0, len(shape.table.columns)):
			print("\nCOLUMN %d HEADER: %s") % i, str(shape.table.cell(0, i).text) 
		"""
		print("\n")


		return True
	return False


def findShapeTEXT_BOX(shape):
	if(string_found('TEXT_BOX',str(shape.shape_type))):
		print("\n FOUND A TEXT_BOX")
		findShapeDimensions(shape)
		"""
		print("\nShape Attributes: ")
		pprint(dir(shape))
		print("\n")
		pprint(dir(shape.text_frame))
		print("\n")
		pprint(dir(shape.text_frame.paragraphs[0]))
		print("\n")
		print(len(shape.text_frame.paragraphs))
		print("\n")
		print(shape.text_frame.paragraphs[0].text)
		"""
		if(shape.text):
			print("\nSHAPE TEXT: " + shape.text)
			return True
		else:
			print("\nSHAPE TEXT: NOT FOUND!!!")
	return False

def findShapeTEXT_EFFECT(shape):
	if(string_found('TEXT_EFFECT',str(shape.shape_type))):
		print("\n FOUND A TEXT_EFFECT")
		return True
	return False

def findShapeWEB_VIDEO(shape):
	if(string_found('WEB_VIDEO',str(shape.shape_type))):
		print("\n FOUND A WEB_VIDEO")
		return True
	return False

def findShapeMIXED(shape):
	if(string_found('MIXED',str(shape.shape_type))):
		print("\n FOUND A MIXED")
		return True
	return False


"""
############################################################################
FIND BY PLACEHOLDER FUNCTIONS
############################################################################
"""

def findPlaceholderTITLE(shape):
	if(string_found('TITLE',str(shape.placeholder_format.type))):
		print("\n FOUND A TITLE")
		findShapeDimensions(shape)
		if(shape.text_frame):
			print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
			return True
		else:
			print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
	return False

def findPlaceholderCENTER_TITLE(shape):
	if(string_found('CENTER_TITLE',str(shape.placeholder_format.type))):
		print("\n FOUND A CENTER_TITLE")
		findShapeDimensions(shape)
		if(shape.text_frame):
			print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
			return True
		else:
			print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
	return False

def findPlaceholderSUBTITLE(shape):
	if(string_found('SUBTITLE',str(shape.placeholder_format.type))):
		print("\n FOUND A SUBTITLE")
		findShapeDimensions(shape)
		if(shape.text_frame):
			print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
			return True
		else:
			print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
	return False

def findPlaceholderTABLE(shape):
	if(string_found('TABLE',str(shape.placeholder_format.type))):
		print("\n FOUND A TABLE")
		return True
	return False

def findPlaceholderPICTURE(shape):
	if(string_found('PICTURE',str(shape.placeholder_format.type))):
		print("\n FOUND A PICTURE")
		return True
	return False

def findPlaceholderCHART(shape):
	if(string_found('CHART',str(shape.placeholder_format.type))):
		print("\n FOUND A CHART")
		return True
	return False

def findPlaceholderBODY(shape):
	if(string_found('BODY',str(shape.placeholder_format.type))):
		print("\n FOUND A BODY")
		findShapeDimensions(shape)
		if(shape.text_frame):
			print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
			return True
		elif(shape.text):
			print("\nSHAPE.TEXT: " + shape.text)
			return True
		else:
			print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
	return False

def findPlaceholderOBJECT(shape):
	if(string_found('OBJECT',str(shape.placeholder_format.type))):
		print("\n FOUND A OBJECT")
		findShapeDimensions(shape)
		if(shape.text_frame):
			print("\nTEXT_FRAME TEXT: " + shape.text_frame.text)
			return True
		else:
			print("\nTEXT_FRAME TEXT: NOT FOUND!!!")
	return False

"""
############################################################################
CORE FUNCTIONS
############################################################################
"""

def findPlaceholder(shape):
	if(string_found('PLACEHOLDER',str(shape.shape_type))):
		print("\n FOUND A PLACEHOLDER ON SHAPE")
		FOUND_BY = 'BY PLACEHOLDER'
		if(findPlaceholderTITLE(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderCENTER_TITLE(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderSUBTITLE(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderTABLE(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderPICTURE(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderCHART(shape)):
			separateFoundShapesWithPrint(FOUND_BY)

		elif(findPlaceholderBODY(shape)):
			separateFoundShapesWithPrint(FOUND_BY)
			
		elif(findPlaceholderOBJECT(shape)):
			separateFoundShapesWithPrint(FOUND_BY)
		else:
			print('\n NO SHAPE IDENTIFIED BY PLACEHOLDER\n')
			return False
		return True
	return False

def findShape(shape):
	print("\n SEARCHING FOR SHAPE BY SHAPE.SHAPE_TYPE: " + str(shape.shape_type))
	FOUND_BY = 'BY SHAPE.SHAPE_TYPE'
	if(findShapeTEXT_BOX(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeAUTO_SHAPE(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeCALLOUT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeCANVAS(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeCHART(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeCOMMENT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeDIAGRAM(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeEMBEDDED_OLE_OBJECT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeFORM_CONTROL(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeFREEFORM(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeGROUP(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeIGX_GRAPHIC(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeINK(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeINK_COMMENT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeLINE(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeLINKED_OLE_OBJECT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeLINKED_PICTURE(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeMEDIA(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeOLE_CONTROL_OBJECT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapePICTURE(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapePLACEHOLDER(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeSCRIPT_ANCHOR(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeTABLE(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeTEXT_BOX(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeTEXT_EFFECT(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeWEB_VIDEO(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	elif(findShapeMIXED(shape)):
		separateFoundShapesWithPrint(FOUND_BY)
	else:
		print("\n NEW SHAPE_TYPE: " + str(shape.shape_type))
		separateFoundShapesWithPrint(FOUND_BY)

def identifyAllShapes(prs):
	for slide in prs.slides:
		print("\n\n#############################################################\n\n")
		print("Slide " + str((prs.slides.index(slide)) + 1) + " has the following shapes:")
		for shape in slide.shapes:
			identifyShape(shape)

def identifyShape(shape):
			try:
				print("\nShape Type: " + str(shape.shape_type))
				
				if not (findPlaceholder(shape)):
					findShape(shape)
			except Exception as detail:
				print(detail)
				pass

def scanPresentationByCLI():
	try: # try if ran at command line with file as first arg
		FILE_NAME = sys.argv[1] # ABSOLUTE FILE PATH OF PPTX TO BE SCANNED
		current_dr = os.path.dirname(os.path.realpath(__file__))
		FULL_PATH = current_dr + "/" + FILE_NAME 
		print(FULL_PATH)
		prs = Presentation(FULL_PATH)
		return prs
	except Exception as detail:
		print(detail)
		pass	

def scanPresentationByMethod(FILE_NAME): # FILE_NAME EXPECTS FILENAME for Scanning (NO QUOTES, OR PREFIX)
	try: # try if ran at command line with file as first arg
		current_dr = os.path.dirname(os.path.realpath(__file__))
		FULL_PATH = current_dr + "/" + FILE_NAME 
		print(FULL_PATH)
		prs = Presentation(FULL_PATH)
		return prs
	except Exception as detail:
		print(detail)
		pass		

"""
############################################################################
MAIN
############################################################################
"""

if __name__ == "__main__":
	prs = scanPresentationByCLI()
	identifyAllShapes(prs)
"""
############################################################################
TEMPLATES
############################################################################
"""

""" def TEMPLATE TYPES TO CHECK PLACEHOLDERS
				USE ALL CAPS
				BITMAP
				    Bitmap
				BODY
				    Body
				CENTER_TITLE
				    Center Title
				CHART
				    Chart
				DATE
				    Date
				FOOTER
				    Footer
				HEADER
				    Header
				MEDIA_CLIP
				    Media Clip
				OBJECT
				    Object
				ORG_CHART
				    Organization Chart
				PICTURE
				    Picture
				SLIDE_NUMBER
				    Slide Number
				SUBTITLE
				    Subtitle
				TABLE
				    Table
				TITLE
				    Title
				VERTICAL_BODY
				    Vertical Body
				VERTICAL_OBJECT
				    Vertical Object
				VERTICAL_TITLE
				    Vertical Title
				MIXED
"""

""" def TEMPLATE TYPES TO CHECK SHAPES
	AUTO_SHAPE

	CALLOUT

	CANVAS

	CHART

	COMMENT

	DIAGRAM

	EMBEDDED_OLE_OBJECT

	FORM_CONTROL

	FREEFORM

	GROUP

	IGX_GRAPHIC

	INK

	INK_COMMENT

	LINE

	LINKED_OLE_OBJECT

	LINKED_PICTURE

	MEDIA

	OLE_CONTROL_OBJECT

	PICTURE

	PLACEHOLDER

	SCRIPT_ANCHOR

	TABLE

	TEXT_BOX

	TEXT_EFFECT

	WEB_VIDEO

	MIXED
"""

""" def TEMPLATE_FUNCTION
def find_text(prs):
	for slide in prs.slides:
		for shape in slide.shapes:
			try:
				if shape.text_frame:
					print("Slide " + str((prs.slides.index(slide)) + 1) + " has text_frame.")
					for paragraph in shape.text_frame.paragraphs:
						text_runs = []
						for run in paragraph.runs:
							text_runs.append(run.text)
							print(run.text)
						print("Slide " + str((prs.slides.index(slide)) + 1) + " has text paragraph: \n" + text_runs)
					#writeTable(__file__)
					continue
			except:
				# print("Slide " + str((prs.slides.index(slide)) + 1) + " has NO table.")
				pass
"""

""" def TEMPLATE_FUNCTION
def findShapeXXX(shape):
	if(
string_found('PLACEHOLDER',str(shape.shape_type))
):
		# CONTENT
		return True
	return False
"""
